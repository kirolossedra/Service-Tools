import requests
from bs4 import BeautifulSoup
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tkinter as tk
from tkinter import scrolledtext, messagebox
import threading
import os

def clean_text(text):
    """Remove ALL special formatting and convert to plain ASCII text."""
    # Convert to ASCII, ignoring anything that can't be converted
    return text.encode('ascii', 'ignore').decode('ascii').strip()


def search_genius_lyrics(song_query):
    """Search for a song on Genius using the API and get the lyrics URL and content."""
    
    # Use the Genius API endpoint
    url = f"https://genius.com/api/search/multi?per_page=5&q={song_query.replace(' ', '%20')}"
    
    print(f"Searching: {url}")
    
    try:
        response = requests.get(url)
        response.raise_for_status()
        data = response.json()
    except requests.RequestException as e:
        print(f"Error fetching search results: {e}")
        return None
    
    # Find the first song result
    lyrics_url = None
    song_title = None
    
    for section in data['response']['sections']:
        if section['type'] == 'song':
            for hit in section['hits']:
                lyrics_url = hit['result']['url']
                # Get title and artist and FORCE them to plain text
                title = clean_text(hit['result']['title'])
                artist = clean_text(hit['result']['primary_artist']['name'])
                song_title = f"{title} by {artist}"
                print(f"Found: {song_title}")
                print(f"URL: {lyrics_url}")
                break
            if lyrics_url:
                break
    
    if not lyrics_url:
        print("No lyrics URL found in search results")
        return None
    
    # Now fetch the lyrics from the URL
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
    }
    
    try:
        lyrics_response = requests.get(lyrics_url, headers=headers)
        lyrics_response.raise_for_status()
    except requests.RequestException as e:
        print(f"Error fetching lyrics page: {e}")
        return None
    
    lyrics_soup = BeautifulSoup(lyrics_response.text, 'html.parser')
    
    # Find lyrics containers
    lyrics_containers = lyrics_soup.find_all('div', attrs={'data-lyrics-container': 'true'})
    
    if not lyrics_containers:
        lyrics_containers = lyrics_soup.find_all('div', class_=re.compile(r'Lyrics__Container'))
    
    if not lyrics_containers:
        print("Could not find lyrics in the page")
        return {'url': lyrics_url, 'title': song_title, 'lyrics': None}
    
    lyrics_text = []
    for container in lyrics_containers:
        for br in container.find_all('br'):
            br.replace_with('\n')
        text = container.get_text()
        lyrics_text.append(text)
    
    lyrics = '\n\n'.join(lyrics_text).strip()
    
    return {
        'url': lyrics_url,
        'title': song_title,
        'lyrics': lyrics
    }


def parse_lyrics_sections(lyrics_text):
    """
    Parse lyrics into sections based on brackets [Section Name].
    Returns list of tuples: (section_name, section_text)
    Only includes sections that have bracket markers.
    """
    sections = []
    
    # Split by sections marked with [...]
    parts = re.split(r'\[([^\]]+)\]', lyrics_text)
    
    # Skip the first part if there's content before any section markers
    for i in range(1, len(parts), 2):
        if i < len(parts):
            section_name = parts[i].strip()
            section_text = parts[i + 1].strip() if i + 1 < len(parts) else ""
            
            if section_text:  # Only add if there's actual content
                sections.append((section_name, section_text))
    
    return sections


def calculate_font_size(text_length):
    """Calculate appropriate font size based on text length with more granular hierarchy."""
    if text_length < 80:
        return 48
    elif text_length < 120:
        return 44
    elif text_length < 180:
        return 40
    elif text_length < 250:
        return 36
    elif text_length < 350:
        return 32
    elif text_length < 450:
        return 28
    elif text_length < 600:
        return 26
    elif text_length < 800:
        return 24
    elif text_length < 1000:
        return 22
    else:
        return 20


def create_lyrics_presentation(song_title, sections, output_file="lyrics_presentation.pptx"):
    """
    Create a PowerPoint presentation with lyrics sections.
    
    Args:
        song_title (str): Title of the song
        sections (list): List of tuples (section_name, section_text)
        output_file (str): Output filename
    """
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title text box
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True  # Enable word wrap
    tf.text = song_title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 255)
    
    print(f"Created title slide: {song_title}")
    
    # Create slides for each section (without the [Section Name] header)
    for section_name, section_text in sections:
        slide = prs.slides.add_slide(blank_layout)
        
        # Calculate font size based on content length
        content_length = len(section_text)
        font_size = calculate_font_size(content_length)
        
        # Add lyrics text (centered, taking full vertical space)
        left = Inches(0.5)
        top = Inches(1)
        width = Inches(9)
        height = Inches(6)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = section_text
        
        # Center align and format
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        print(f"Created slide: [{section_name}] ({content_length} chars, {font_size}pt font)")
    
    # Save presentation
    prs.save(output_file)
    print(f"\n✓ Presentation saved as: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")


def process_single_song(song_query, output_folder, log_callback):
    """Process a single song and create presentation."""
    try:
        log_callback(f"\n{'='*60}")
        log_callback(f"Processing: {song_query}")
        
        # Search and get lyrics
        result = search_genius_lyrics(song_query)
        
        if not result or not result['lyrics']:
            log_callback(f"❌ Could not retrieve lyrics for: {song_query}")
            return False
        
        log_callback(f"✓ Found: {result['title']}")
        
        # Parse lyrics into sections
        sections = parse_lyrics_sections(result['lyrics'])
        log_callback(f"✓ Parsed {len(sections)} sections")
        
        # Create safe filename
        safe_filename = "".join(c for c in song_query if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_filename = safe_filename.replace(' ', '_')
        output_file = os.path.join(output_folder, f"{safe_filename}.pptx")
        
        # Create PowerPoint
        create_lyrics_presentation(result['title'], sections, output_file)
        log_callback(f"✓ Saved: {output_file}")
        
        return True
    except Exception as e:
        log_callback(f"❌ Error processing '{song_query}': {str(e)}")
        return False


class LyricsGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("Lyrics Presentation Generator")
        self.root.geometry("700x600")
        
        # Title
        title_label = tk.Label(root, text="Lyrics Presentation Generator", 
                              font=("Arial", 16, "bold"))
        title_label.pack(pady=10)
        
        # Instructions
        instructions = tk.Label(root, text="Enter song queries (one per line):", 
                               font=("Arial", 10))
        instructions.pack(pady=5)
        
        # Text input area
        self.text_area = scrolledtext.ScrolledText(root, width=80, height=15, 
                                                   font=("Arial", 10))
        self.text_area.pack(pady=10, padx=20)
        self.text_area.insert(tk.END, "king of my heart bethel\noceans hillsong\ngoodness of god")
        
        # Output folder
        folder_frame = tk.Frame(root)
        folder_frame.pack(pady=5)
        
        tk.Label(folder_frame, text="Output Folder:", font=("Arial", 10)).pack(side=tk.LEFT, padx=5)
        self.folder_entry = tk.Entry(folder_frame, width=40, font=("Arial", 10))
        self.folder_entry.insert(0, "lyrics_presentations")
        self.folder_entry.pack(side=tk.LEFT, padx=5)
        
        # Generate button
        self.generate_button = tk.Button(root, text="Generate Presentations", 
                                        command=self.start_processing,
                                        font=("Arial", 12, "bold"),
                                        bg="#4CAF50", fg="white",
                                        padx=20, pady=10)
        self.generate_button.pack(pady=10)
        
        # Log area
        log_label = tk.Label(root, text="Progress Log:", font=("Arial", 10))
        log_label.pack(pady=5)
        
        self.log_area = scrolledtext.ScrolledText(root, width=80, height=10, 
                                                  font=("Courier", 9),
                                                  bg="#f0f0f0")
        self.log_area.pack(pady=5, padx=20)
        
    def log(self, message):
        """Add message to log area."""
        self.log_area.insert(tk.END, message + "\n")
        self.log_area.see(tk.END)
        self.root.update_idletasks()
    
    def start_processing(self):
        """Start processing songs in a separate thread."""
        # Get input
        text = self.text_area.get("1.0", tk.END).strip()
        if not text:
            messagebox.showwarning("No Input", "Please enter at least one song query!")
            return
        
        songs = [line.strip() for line in text.split('\n') if line.strip()]
        output_folder = self.folder_entry.get().strip()
        
        if not output_folder:
            output_folder = "lyrics_presentations"
        
        # Disable button during processing
        self.generate_button.config(state=tk.DISABLED, text="Processing...")
        self.log_area.delete("1.0", tk.END)
        
        # Process in thread
        thread = threading.Thread(target=self.process_songs, args=(songs, output_folder))
        thread.daemon = True
        thread.start()
    
    def process_songs(self, songs, output_folder):
        """Process all songs."""
        try:
            # Create output folder
            if not os.path.exists(output_folder):
                os.makedirs(output_folder)
                self.log(f"✓ Created output folder: {output_folder}")
            
            self.log(f"\nProcessing {len(songs)} song(s)...\n")
            
            success_count = 0
            fail_count = 0
            
            for song in songs:
                if process_single_song(song, output_folder, self.log):
                    success_count += 1
                else:
                    fail_count += 1
            
            self.log(f"\n{'='*60}")
            self.log(f"COMPLETED!")
            self.log(f"✓ Success: {success_count}")
            self.log(f"❌ Failed: {fail_count}")
            self.log(f"Output folder: {os.path.abspath(output_folder)}")
            
            messagebox.showinfo("Complete", 
                              f"Processing complete!\n\n"
                              f"Success: {success_count}\n"
                              f"Failed: {fail_count}\n\n"
                              f"Files saved in: {output_folder}")
            
        except Exception as e:
            self.log(f"\n❌ ERROR: {str(e)}")
            messagebox.showerror("Error", f"An error occurred:\n{str(e)}")
        
        finally:
            # Re-enable button
            self.generate_button.config(state=tk.NORMAL, text="Generate Presentations")


def main():
    """Launch the GUI application."""
    root = tk.Tk()
    app = LyricsGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()
