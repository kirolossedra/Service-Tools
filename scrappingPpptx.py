import requests
from bs4 import BeautifulSoup
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def search_genius_lyrics(song_query):
    """Search for a song on Genius using the API and get the lyrics URL and content."""
    
    # Use the Genius API endpoint
    url = f"https://genius.com/api/search/multi?per_page=5&q={song_query.replace(' ', '%20')}"
    
    print(f"Searching: {url}")
    
    try:
        response = requests.get(url)
        response.raise_for_status()
        data = response.json()
    except requests.RequestException as e:
        print(f"Error fetching search results: {e}")
        return None
    
    # Find the first song result
    lyrics_url = None
    song_title = None
    
    for section in data['response']['sections']:
        if section['type'] == 'song':
            for hit in section['hits']:
                lyrics_url = hit['result']['url']
                song_title = hit['result']['full_title']
                print(f"Found: {song_title}")
                print(f"URL: {lyrics_url}")
                break
            if lyrics_url:
                break
    
    if not lyrics_url:
        print("No lyrics URL found in search results")
        return None
    
    # Now fetch the lyrics from the URL
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
    }
    
    try:
        lyrics_response = requests.get(lyrics_url, headers=headers)
        lyrics_response.raise_for_status()
    except requests.RequestException as e:
        print(f"Error fetching lyrics page: {e}")
        return None
    
    lyrics_soup = BeautifulSoup(lyrics_response.text, 'html.parser')
    
    # Find lyrics containers
    lyrics_containers = lyrics_soup.find_all('div', attrs={'data-lyrics-container': 'true'})
    
    if not lyrics_containers:
        lyrics_containers = lyrics_soup.find_all('div', class_=re.compile(r'Lyrics__Container'))
    
    if not lyrics_containers:
        print("Could not find lyrics in the page")
        return {'url': lyrics_url, 'title': song_title, 'lyrics': None}
    
    lyrics_text = []
    for container in lyrics_containers:
        for br in container.find_all('br'):
            br.replace_with('\n')
        text = container.get_text()
        lyrics_text.append(text)
    
    lyrics = '\n\n'.join(lyrics_text).strip()
    
    return {
        'url': lyrics_url,
        'title': song_title,
        'lyrics': lyrics
    }


def parse_lyrics_sections(lyrics_text):
    """
    Parse lyrics into sections based on brackets [Section Name].
    Returns list of tuples: (section_name, section_text)
    """
    sections = []
    
    # Split by sections marked with [...]
    parts = re.split(r'\[([^\]]+)\]', lyrics_text)
    
    current_section = None
    
    for i, part in enumerate(parts):
        part = part.strip()
        if not part:
            continue
            
        # Odd indices are section names (from the capturing group)
        if i % 2 == 1:
            current_section = part
        # Even indices are the content following the section
        else:
            if current_section:
                sections.append((current_section, part))
                current_section = None
            elif part:  # Content without a section header
                sections.append(("", part))
    
    return sections


def calculate_font_size(text_length):
    """Calculate appropriate font size based on text length."""
    if text_length < 100:
        return 44
    elif text_length < 200:
        return 36
    elif text_length < 300:
        return 32
    elif text_length < 500:
        return 28
    elif text_length < 700:
        return 24
    else:
        return 20


def create_lyrics_presentation(song_title, sections, output_file="lyrics_presentation.pptx"):
    """
    Create a PowerPoint presentation with lyrics sections.
    
    Args:
        song_title (str): Title of the song
        sections (list): List of tuples (section_name, section_text)
        output_file (str): Output filename
    """
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title text box
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = song_title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 255)
    
    print(f"Created title slide: {song_title}")
    
    # Create slides for each section
    for section_name, section_text in sections:
        slide = prs.slides.add_slide(blank_layout)
        
        # Calculate font size based on content length
        content_length = len(section_text)
        font_size = calculate_font_size(content_length)
        
        # Add section name if it exists (at top)
        if section_name:
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(9)
            height = Inches(0.8)
            
            section_box = slide.shapes.add_textbox(left, top, width, height)
            section_tf = section_box.text_frame
            section_tf.text = f"[{section_name}]"
            section_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            section_tf.paragraphs[0].font.size = Pt(24)
            section_tf.paragraphs[0].font.bold = True
            section_tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 150)
        
        # Add lyrics text (centered)
        left = Inches(0.5)
        top = Inches(1.5) if section_name else Inches(1)
        width = Inches(9)
        height = Inches(5.5) if section_name else Inches(6)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = section_text
        
        # Center align and format
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        print(f"Created slide: [{section_name}] ({content_length} chars, {font_size}pt font)")
    
    # Save presentation
    prs.save(output_file)
    print(f"\n✓ Presentation saved as: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")


def main(song_query, output_file="lyrics_presentation.pptx"):
    """Main function to search, parse, and create presentation."""
    
    # Step 1: Search and get lyrics
    print(f"Searching for: {song_query}\n")
    result = search_genius_lyrics(song_query)
    
    if not result or not result['lyrics']:
        print("Could not retrieve lyrics")
        return
    
    print(f"\nFound: {result['title']}")
    print(f"URL: {result['url']}\n")
    
    # Step 2: Parse lyrics into sections
    sections = parse_lyrics_sections(result['lyrics'])
    print(f"Parsed {len(sections)} sections\n")
    
    # Step 3: Create PowerPoint
    print("Creating PowerPoint presentation...\n")
    create_lyrics_presentation(result['title'], sections, output_file)


# Example usage
if __name__ == "__main__":
    # You can change the song query here
    song_query = "king of my heart bethel"
    output_file = "king_of_my_heart_lyrics.pptx"
    
    main(song_query, output_file)
    
    print("\n" + "="*60)
    print("To use with a different song:")
    print('  main("oceans hillsong", "oceans_lyrics.pptx")')
    print('  main("goodness of god", "goodness_lyrics.pptx")')
